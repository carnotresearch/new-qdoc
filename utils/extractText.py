"""
Simplified text extraction module for various document formats with page tracking.

This module provides unified text extraction with automatic file type detection
and consistent output formatting across all supported file types.
"""

import asyncio
import logging
import mimetypes
import os
import re
import tempfile
import time
import unicodedata
from typing import Dict, List, Tuple, Any

# Third-party imports
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from langchain.schema import Document
from openpyxl import load_workbook
from pptx import Presentation
from werkzeug.datastructures import FileStorage
from .extractTextOCR import extract_pdf_ocr
from config import Config


# Configure logging
logger = logging.getLogger(__name__)

def clean_filename(filename: str) -> str:
    """Clean filename by removing temporary path prefixes"""
    if not filename:
        return filename
    
    # Handle both Windows and Unix paths
    if '\\' in filename:
        filename = filename.split('\\')[-1]
    if '/' in filename:
        filename = filename.split('/')[-1]
    
    return filename

def clean_text(text: str) -> str:
    """Apply universal text cleaning across all file types"""
    if not text:
        return text
    
    # Normalize Unicode characters
    text = unicodedata.normalize("NFKC", text)
    
    # Replace common problematic characters
    replacements = {
        "ﬁ": "fi", "ﬂ": "fl", "ﬀ": "ff", "ﬃ": "ffi", "ﬄ": "ffl",
        """: '"', """: '"', "'": "'", "'": "'", "′": "'",
        "‒": "-", "–": "-", "—": "-", "―": "-",
        "…": "...", "•": "*", "°": " degrees ",
        "©": "(c)", "®": "(R)", "™": "(TM)"
    }
    
    for old, new in replacements.items():
        text = text.replace(old, new)
    
    # Clean control characters while preserving whitespace
    text = "".join(
        char for char in text
        if unicodedata.category(char)[0] != "C" or char in "\n\t "
    )
    
    # Clean up spacing
    text = re.sub(r"[ \t]+", " ", text)  # Consolidate horizontal whitespace
    text = re.sub(r" +\n", "\n", text)  # Remove spaces before newlines
    text = re.sub(r"\n +", "\n", text)  # Remove spaces after newlines
    text = re.sub(r"\n{3,}", "\n\n", text)  # Max two consecutive newlines
    text = re.sub(r"^\s+", "", text)  # Remove leading whitespace
    text = re.sub(r"\s+$", "", text)  # Remove trailing whitespace
    
    # Clean up punctuation spacing
    text = re.sub(r"\s+([.,;:!?)])", r"\1", text)
    text = re.sub(r"(\()\s+", r"\1", text)
    
    # Remove zero-width characters
    text = re.sub(r"[\u200b\u200c\u200d\ufeff\u200e\u200f]", "", text)
    
    # Fix hyphenation at line breaks
    text = re.sub(r"(?<=\w)-\s*\n\s*(?=\w)", "", text)
    
    return text.strip()

def detect_file_type(file_path: str, filename: str = None) -> str:
    """Detect file type using MIME types and extensions"""
    # Try MIME type detection first
    mime_type, _ = mimetypes.guess_type(file_path)
    if mime_type:
        return mime_type
    
    # Fallback to extension
    ext = os.path.splitext(filename or file_path)[1].lower()
    
    ext_mime_map = {
        ".txt": "text/plain",
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    }
    
    return ext_mime_map.get(ext, "text/plain")

async def extract_txt(file_path: str, filename: str) -> List[Document]:
    """Extract text from plain text files - treat as single page"""
    def _read_file():
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()
    
    content = await asyncio.get_event_loop().run_in_executor(None, _read_file)
    cleaned_content = clean_text(content)
    
    if not cleaned_content.strip():
        return []
    
    # Create a single Document for the text file
    doc = Document(
        page_content=cleaned_content,
        metadata={
            "source": filename,
            "filename": filename,
            "page": 0  # Text files are treated as single page (0-indexed)
        }
    )
    
    return [doc]

# async def extract_pdf(file_path: str, filename: str) -> List[Document]:
    """Extract text from PDF files with enhanced Unicode/Hindi support"""
    def _extract():
        doc = fitz.open(file_path)
        try:
            pages = []
            for page_num, page in enumerate(doc):
                # Try multiple extraction methods for better Unicode handling
                text = _extract_text_with_fallback(page)
                
                if text and text.strip():  # Only include pages with content
                    pages.append({
                        "content": text,
                        "page_number": page_num  # 0-indexed
                    })
            return pages
        finally:
            doc.close()
    
    pages_data = await asyncio.get_event_loop().run_in_executor(None, _extract)
    
    documents = []
    for page_data in pages_data:
        cleaned_content = clean_text_unicode(page_data["content"])
        if cleaned_content.strip():
            doc = Document(
                page_content=cleaned_content,
                metadata={
                    "source": filename,
                    "filename": filename,
                    "page": page_data["page_number"]
                }
            )
            documents.append(doc)
    
    return documents


async def extract_pdf(file_path: str, filename: str) -> List[Document]:
    """Extract text from PDF files - one Document per page"""
    def _extract():
        doc = fitz.open(file_path)
        try:
            pages = []
            for page_num, page in enumerate(doc):
                text = page.get_text()
                if text.strip():  # Only include pages with content
                    pages.append({
                        "content": text,
                        "page_number": page_num  # 0-indexed
                    })
            return pages
        finally:
            doc.close()
    
    pages_data = await asyncio.get_event_loop().run_in_executor(None, _extract)
    
    documents = []
    for page_data in pages_data:
        cleaned_content = clean_text(page_data["content"])
        if cleaned_content.strip():
            doc = Document(
                page_content=cleaned_content,
                metadata={
                    "source": filename,
                    "filename": filename,
                    "page": page_data["page_number"]
                }
            )
            documents.append(doc)
    
    return documents


def _extract_text_with_fallback(page):
    """
    Extract text using multiple methods with fallback for better Unicode support
    """
    try:
        # Method 1: Default text extraction
        text = page.get_text()
        if _is_readable_text(text):
            logger.debug("Using default text extraction")
            return text
        
        # Method 2: Extract with layout preservation
        text = page.get_text("text", flags=fitz.TEXT_PRESERVE_LIGATURES | fitz.TEXT_PRESERVE_WHITESPACE)
        if _is_readable_text(text):
            logger.debug("Using layout-preserved text extraction")
            return text
            
        # Method 3: Extract with different encoding handling
        text = page.get_text("rawdict")
        if text:
            extracted_text = _extract_from_rawdict(text)
            if _is_readable_text(extracted_text):
                logger.debug("Using rawdict text extraction")
                return extracted_text
        
        # Method 4: Try HTML extraction and clean it
        html_text = page.get_text("html")
        if html_text:
            import re
            # Remove HTML tags and extract text
            clean_html = re.sub(r'<[^>]+>', '', html_text)
            if _is_readable_text(clean_html):
                logger.debug("Using HTML text extraction")
                return clean_html
        
        # Method 5: Extract using text blocks with position info
        blocks = page.get_text("dict")
        if blocks:
            extracted_text = _extract_from_blocks(blocks)
            if _is_readable_text(extracted_text):
                logger.debug("Using blocks text extraction")
                return extracted_text
        
        # If all methods fail, return the original text anyway
        logger.warning(f"All extraction methods failed for page, using default")
        return page.get_text()
        
    except Exception as e:
        logger.error(f"Error in text extraction: {e}")
        # Fallback to basic extraction
        return page.get_text()

def _is_readable_text(text):
    """
    Check if extracted text appears to be readable (not gibberish)
    """
    if not text or len(text.strip()) < 10:
        return False
    
    # Check for reasonable ratio of alphanumeric/unicode characters
    total_chars = len(text)
    readable_chars = sum(1 for c in text if c.isalnum() or c.isspace() or unicodedata.category(c).startswith('L'))
    
    # Should have at least 60% readable characters
    readable_ratio = readable_chars / total_chars if total_chars > 0 else 0
    
    # Also check for excessive punctuation which might indicate encoding issues
    punct_chars = sum(1 for c in text if unicodedata.category(c).startswith('P'))
    punct_ratio = punct_chars / total_chars if total_chars > 0 else 0
    
    return readable_ratio > 0.6 and punct_ratio < 0.3

def _extract_from_rawdict(rawdict_data):
    """Extract text from PyMuPDF rawdict format"""
    try:
        extracted_text = []
        if 'blocks' in rawdict_data:
            for block in rawdict_data['blocks']:
                if 'lines' in block:
                    for line in block['lines']:
                        if 'spans' in line:
                            line_text = []
                            for span in line['spans']:
                                if 'text' in span:
                                    # Try to fix encoding issues
                                    span_text = _fix_encoding(span['text'])
                                    line_text.append(span_text)
                            if line_text:
                                extracted_text.append(' '.join(line_text))
        return '\n'.join(extracted_text)
    except Exception as e:
        logger.error(f"Error extracting from rawdict: {e}")
        return ""

def _extract_from_blocks(blocks_data):
    """Extract text from PyMuPDF blocks format"""
    try:
        extracted_text = []
        if 'blocks' in blocks_data:
            for block in blocks_data['blocks']:
                if block.get('type') == 0:  # Text block
                    if 'lines' in block:
                        for line in block['lines']:
                            if 'spans' in line:
                                line_text = []
                                for span in line['spans']:
                                    if 'text' in span:
                                        span_text = _fix_encoding(span['text'])
                                        line_text.append(span_text)
                                if line_text:
                                    extracted_text.append(' '.join(line_text))
        return '\n'.join(extracted_text)
    except Exception as e:
        logger.error(f"Error extracting from blocks: {e}")
        return ""

def _fix_encoding(text):
    """
    Attempt to fix common encoding issues in extracted text
    """
    if not text:
        return text
    
    try:
        # First, ensure it's properly decoded as UTF-8
        if isinstance(text, bytes):
            text = text.decode('utf-8', errors='ignore')
        
        # Try to fix common Hindi encoding issues
        # Some PDFs use incorrect encodings for Devanagari script
        
        # Method 1: Try to re-encode and decode with different encodings
        encodings_to_try = ['utf-8', 'utf-16', 'iso-8859-1', 'cp1252', 'latin1']
        
        for encoding in encodings_to_try:
            try:
                if isinstance(text, str):
                    # Convert to bytes then back to string with different encoding
                    test_bytes = text.encode('latin1', errors='ignore')
                    decoded_text = test_bytes.decode(encoding, errors='ignore')
                    
                    # Check if this produces better results for Hindi
                    if _contains_proper_devanagari(decoded_text):
                        return decoded_text
            except (UnicodeDecodeError, UnicodeEncodeError):
                continue
        
        # Method 2: Fix specific character mappings common in Hindi PDFs
        char_fixes = {
            'ﬁ': 'फि',
            'ﬂ': 'फ्ल',
            '¼': '्र',
            '½': '्य',
            '¾': '्व',
            'ò': 'ो',
            'è': 'े',
            'ì': 'ि',
            'î': 'ी',
            'ù': 'ु',
            'û': 'ू',
            '`': 'ं',
            '~': 'ँ',
        }
        
        fixed_text = text
        for wrong, correct in char_fixes.items():
            fixed_text = fixed_text.replace(wrong, correct)
        
        return fixed_text
        
    except Exception as e:
        logger.error(f"Error fixing encoding: {e}")
        return text

def _contains_proper_devanagari(text):
    """Check if text contains proper Devanagari characters"""
    if not text:
        return False
    
    # Count Devanagari characters (Unicode range 0x0900-0x097F)
    devanagari_count = sum(1 for c in text if '\u0900' <= c <= '\u097F')
    total_chars = len([c for c in text if not c.isspace()])
    
    # If more than 10% of non-space characters are Devanagari, consider it proper
    return devanagari_count > 0 and (devanagari_count / max(total_chars, 1)) > 0.1

def clean_text_unicode(text: str) -> str:
    """Enhanced text cleaning with better Unicode support"""
    if not text:
        return text
    
    # First apply the original cleaning
    text = clean_text(text)
    
    # Additional Unicode normalization for Hindi/Devanagari
    try:
        # Normalize Unicode to handle different representations of the same character
        text = unicodedata.normalize("NFC", text)
        
        # Remove any remaining problematic control characters but preserve Devanagari
        cleaned_chars = []
        for char in text:
            cat = unicodedata.category(char)
            # Keep letters, numbers, spaces, punctuation, and marks (for diacritics)
            if (cat.startswith(('L', 'N', 'Z', 'P', 'M')) or 
                char in '\n\t ' or 
                '\u0900' <= char <= '\u097F'):  # Devanagari range
                cleaned_chars.append(char)
        
        text = ''.join(cleaned_chars)
        
        return text.strip()
        
    except Exception as e:
        logger.error(f"Error in Unicode cleaning: {e}")
        return text

async def extract_docx(file_path: str, filename: str) -> List[Document]:
    """Extract content from DOCX files - try to detect page breaks or treat as single page"""
    def _extract():
        doc = DocxDocument(file_path)
        content = []
        current_page_content = []
        current_page = 0
        
        for paragraph in doc.paragraphs:
            if not paragraph.text.strip():
                continue
            
            style = paragraph.style.name if paragraph.style else "Normal"
            text = paragraph.text.strip()
            
            # Check for page break
            if paragraph._element.xpath('.//w:br[@w:type="page"]'):
                # Page break found - save current page and start new one
                if current_page_content:
                    content.append({
                        "content": "\n\n".join(current_page_content),
                        "page_number": current_page
                    })
                    current_page_content = []
                    current_page += 1
            
            # Handle headings
            if "Heading" in style:
                level = style[-1] if style[-1].isdigit() else "1"
                heading_marks = "#" * int(level)
                current_page_content.append(f"\n{heading_marks} {text}\n")
            else:
                # Handle text formatting
                formatted_text = []
                for run in paragraph.runs:
                    if run.bold:
                        formatted_text.append(f"**{run.text}**")
                    elif run.italic:
                        formatted_text.append(f"*{run.text}*")
                    else:
                        formatted_text.append(run.text)
                current_page_content.append(''.join(formatted_text))
        
        # Add the last page
        if current_page_content:
            content.append({
                "content": "\n\n".join(current_page_content),
                "page_number": current_page
            })
        
        # If no page breaks were found, treat as single page
        if not content and current_page_content:
            content.append({
                "content": "\n\n".join(current_page_content),
                "page_number": 0
            })
        
        return content
    
    pages_data = await asyncio.get_event_loop().run_in_executor(None, _extract)
    
    documents = []
    for page_data in pages_data:
        cleaned_content = clean_text(page_data["content"])
        if cleaned_content.strip():
            doc = Document(
                page_content=cleaned_content,
                metadata={
                    "source": filename,
                    "filename": filename,
                    "page": page_data["page_number"]
                }
            )
            documents.append(doc)
    
    return documents

async def extract_pptx(file_path: str, filename: str) -> List[Document]:
    """Extract content from PPTX files - one Document per slide"""
    def _extract():
        prs = Presentation(file_path)
        slides_data = []
        
        for slide_number, slide in enumerate(prs.slides):
            content = []
            content.append(f"\n# Slide {slide_number + 1}\n")
            
            # Extract title
            if slide.shapes.title and slide.shapes.title.text.strip():
                content.append(f"## {slide.shapes.title.text.strip()}\n")
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape != slide.shapes.title:
                        content.append(shape.text.strip())
            
            slides_data.append({
                "content": "\n\n".join(content),
                "page_number": slide_number  # 0-indexed
            })
        
        return slides_data
    
    slides_data = await asyncio.get_event_loop().run_in_executor(None, _extract)
    
    documents = []
    for slide_data in slides_data:
        cleaned_content = clean_text(slide_data["content"])
        if cleaned_content.strip():
            doc = Document(
                page_content=cleaned_content,
                metadata={
                    "source": filename,
                    "filename": filename,
                    "page": slide_data["page_number"]
                }
            )
            documents.append(doc)
    
    return documents

async def extract_xlsx(file_path: str, filename: str) -> List[Document]:
    """Extract content from XLSX files - one Document per sheet"""
    def _extract():
        wb = load_workbook(file_path, data_only=True)
        sheets_data = []
        
        for sheet_index, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]
            content = []
            content.append(f"\n# Sheet: {sheet_name}\n")
            
            max_row = min(ws.max_row or 1, 1000)
            max_col = min(ws.max_column or 1, 50)
            
            if max_row <= 1:
                content.append("(Empty sheet)")
            else:
                # Create table
                headers = []
                for col in range(1, max_col + 1):
                    cell_value = ws.cell(row=1, column=col).value
                    headers.append(str(cell_value) if cell_value is not None else "")
                
                content.append("| " + " | ".join(headers) + " |")
                content.append("| " + " | ".join(["---"] * len(headers)) + " |")
                
                # Add data rows
                for row in range(2, max_row + 1):
                    row_data = []
                    for col in range(1, max_col + 1):
                        cell_value = ws.cell(row=row, column=col).value
                        row_data.append(str(cell_value) if cell_value is not None else "")
                    content.append("| " + " | ".join(row_data) + " |")
            
            sheets_data.append({
                "content": "\n".join(content),
                "page_number": sheet_index  # Sheet index as page number
            })
        
        return sheets_data
    
    sheets_data = await asyncio.get_event_loop().run_in_executor(None, _extract)
    
    documents = []
    for sheet_data in sheets_data:
        cleaned_content = clean_text(sheet_data["content"])
        if cleaned_content.strip():
            doc = Document(
                page_content=cleaned_content,
                metadata={
                    "source": filename,
                    "filename": filename,
                    "page": sheet_data["page_number"]
                }
            )
            documents.append(doc)
    
    return documents

async def extract_content_from_file(file_path: str, filename: str = None) -> Dict[str, Any]:
    """
    Extract content from any supported file type with page tracking
    
    Args:
        file_path: Path to the file
        filename: Original filename
        
    Returns:
        Dictionary with extracted documents and metadata
    """
    clean_name = clean_filename(filename or os.path.basename(file_path))
    
    # Validate file
    if not os.path.exists(file_path) or os.path.getsize(file_path) == 0:
        return {
            "documents": [],
            "metadata": {"filename": clean_name, "error": "File not found or empty"},
            "success": False
        }
    
    # Detect file type
    file_type = detect_file_type(file_path, filename)
    
    try:
        # Extract based on file type
        if file_type == "text/plain":
            documents = await extract_txt(file_path, clean_name)
        elif file_type == "application/pdf":
            logging.info(f"OCR INITIALIZATION")
            documents = await extract_pdf_ocr(file_path,filename) 
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            documents = await extract_docx(file_path, clean_name)
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            documents = await extract_pptx(file_path, clean_name)
        elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            documents = await extract_xlsx(file_path, clean_name)
        else:
            # Fallback to text
            documents = await extract_txt(file_path, clean_name)
        
        if not documents:
            return {
                "documents": [],
                "metadata": {"filename": clean_name, "error": "No content extracted"},
                "success": False
            }
        
        total_content_length = sum(len(doc.page_content) for doc in documents)
        
        return {
            "documents": documents,
            "metadata": {
                "filename": clean_name,
                "file_type": file_type,
                "total_pages": len(documents),
                "content_length": total_content_length
            },
            "success": True
        }
        
    except Exception as e:
        logger.error(f"Error extracting content from {clean_name}: {e}")
        return {
            "documents": [],
            "metadata": {"filename": clean_name, "error": str(e)},
            "success": False
        }

def get_text_from_files(files: List[FileStorage]) -> Tuple[List[Document], List[Dict[str, Any]]]:
    """
    Extract text from multiple files (main interface for document_service.py)
    
    Args:
        files: List of file objects
        
    Returns:
        Tuple of (list of Document objects, list of file_info dictionaries)
    """
    async def _async_extract():
        all_documents = []
        file_infos = []
        
        for file in files:
            temp_path = None
            try:
                # Create temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.filename)[1]) as temp_file:
                    file.save(temp_file.name)
                    temp_path = temp_file.name
                
                # Extract content
                result = await extract_content_from_file(temp_path, file.filename)
                
                if result["success"] and result["documents"]:
                    logger.info(f"Successfully extracted {len(result['documents'])} pages from {file.filename}")
                    # Add all documents from this file
                    all_documents.extend(result["documents"])
                
                # File info for tracking
                file_info = {
                    "filename": file.filename,
                    "success": result["success"],
                    "page_count": len(result["documents"]) if result["documents"] else 0,
                    "content_length": result["metadata"].get("content_length", 0)
                }
                
                if not result["success"]:
                    file_info["error"] = result["metadata"].get("error", "Unknown error")
                
                file_infos.append(file_info)
                
            except Exception as e:
                logger.error(f"Error processing file {file.filename}: {e}")
                file_infos.append({
                    "filename": file.filename,
                    "success": False,
                    "error": str(e)
                })
            finally:
                # Clean up temporary file
                if temp_path and os.path.exists(temp_path):
                    os.unlink(temp_path)
        
        return all_documents, file_infos
    
    # Run async function synchronously
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    try:
        return loop.run_until_complete(_async_extract())
    finally:
        loop.close()

def process_single_file(file: FileStorage) -> Tuple[List[Document], Dict[str, Any]]:
    """Legacy function for single file processing"""
    docs, file_infos = get_text_from_files([file])
    return docs, file_infos[0] if file_infos else {"success": False, "error": "No file info"}

# Standalone usage support
async def extract_from_path(file_path: str) -> Dict[str, Any]:
    """Extract content from a file path (for standalone usage)"""
    if not os.path.exists(file_path):
        return {"success": False, "error": f"File not found: {file_path}"}
    
    return await extract_content_from_file(file_path)

def main():
    """Main function for standalone usage"""
    import argparse
    import sys
    
    parser = argparse.ArgumentParser(description="Extract text from documents")
    parser.add_argument("--file", "-f", help="Path to file")
    parser.add_argument("--output", "-o", help="Output file path")
    
    args = parser.parse_args()
    
    if not args.file:
        parser.print_help()
        sys.exit(1)
    
    async def _extract():
        result = await extract_from_path(args.file)
        
        if result["success"]:
            documents = result["documents"]
            print(f"Successfully extracted text from: {result['metadata']['filename']}")
            print(f"Total pages: {len(documents)}")
            print(f"Total content length: {result['metadata']['content_length']} characters")
            print("-" * 50)
            
            if args.output:
                with open(args.output, 'w', encoding='utf-8') as f:
                    for i, doc in enumerate(documents):
                        f.write(f"=== Page {doc.metadata['page'] + 1} ===\n")
                        f.write(doc.page_content)
                        f.write(f"\n\n")
                print(f"Content saved to: {args.output}")
            else:
                print("Extracted content (first page preview):")
                if documents:
                    content = documents[0].page_content
                    print(content[:1000] + "..." if len(content) > 1000 else content)
        else:
            print(f"Error: {result['metadata']['error']}")
            sys.exit(1)
    
    asyncio.run(_extract())

if __name__ == "__main__":
    main()